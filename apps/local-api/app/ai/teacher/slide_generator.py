"""PowerPoint Presentation (.pptx) Generator for Lecture Slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from typing import Dict, Any, List
from app.config import settings
from app.rag.retrieval.retriever import HybridRetriever


class LectureSlideGenerator:
    def __init__(self, subject_id: str):
        self.subject_id = subject_id
        self.subject_dir = settings.SUBJECTS_DIR / subject_id
        self.output_dir = self.subject_dir / "generated_materials"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.retriever = HybridRetriever(subject_id=subject_id)

    async def generate_presentation(
        self,
        topic: str,
        unit_title: str = "Unit 3: Machine Learning Foundations",
        target_slides: int = 5
    ) -> Dict[str, Any]:
        """Generates and writes a real .pptx presentation file using curriculum knowledge."""
        # 1. Retrieve curriculum context for topic
        retrieved = await self.retriever.retrieve(query=topic, limit=3)
        context_snippets = [r["text_content"] for r in retrieved] if retrieved else [
            "Mathematical formulation and objective definitions.",
            "Key algorithmic properties and trade-offs.",
            "Practical implementations and real-world evaluation."
        ]

        # 2. Build Slide Content Outline
        slides_data = [
            {
                "title": topic,
                "subtitle": f"{self.subject_id.replace('-', ' ').title()} — {unit_title}",
                "type": "title"
            },
            {
                "title": "1. Introduction & Motivation",
                "bullets": [
                    f"Core mathematical definition of {topic}.",
                    "Why this technique is required in standard curriculum workflows.",
                    "Relation to empirical risk minimization and model generalization."
                ],
                "type": "content"
            },
            {
                "title": "2. Theoretical Framework & Mechanism",
                "bullets": [
                    "Formulation of the primary objective and loss penalty terms.",
                    f"Context Insight: {context_snippets[0][:130]}...",
                    "Geometric interpretation and parameter coefficient constraints."
                ],
                "type": "content"
            },
            {
                "title": "3. Comparison & Practical Considerations",
                "bullets": [
                    "Computational complexity and gradient updates.",
                    "Hyperparameter tuning via cross-validation techniques.",
                    "Mitigating overfitting in high-dimensional feature spaces."
                ],
                "type": "content"
            },
            {
                "title": "4. Summary & University Exam Focus",
                "bullets": [
                    "Key formulas to remember for 5-mark and 10-mark questions.",
                    "Standard PYQ derivations associated with this topic.",
                    "Review questions and recommended textbook exercises."
                ],
                "type": "content"
            }
        ]

        # 3. Create PPTX via python-pptx
        prs = Presentation()
        # Set 16:9 widescreen dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for s_info in slides_data:
            if s_info["type"] == "title":
                # Title slide layout (layout 0)
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = s_info["title"]
                subtitle.text = s_info["subtitle"]
            else:
                # Title & Content layout (layout 1)
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = s_info["title"]
                
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                
                bullets = s_info.get("bullets", [])
                if bullets:
                    tf.text = bullets[0]
                    for b in bullets[1:]:
                        p = tf.add_paragraph()
                        p.text = b
                        p.level = 0

        # Save to disk
        clean_slug = "".join(c if c.isalnum() else "_" for c in topic).strip("_")
        filename = f"{clean_slug}_Lecture_Slides.pptx"
        output_file = self.output_dir / filename
        prs.save(str(output_file))

        return {
            "status": "generated",
            "topic": topic,
            "filename": filename,
            "file_path": str(output_file),
            "size_bytes": output_file.stat().st_size,
            "slides_count": len(slides_data),
            "slides": slides_data
        }
